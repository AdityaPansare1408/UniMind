from pathlib import Path

from pptx import Presentation

from langchain_core.documents import Document as LangChainDocument

from src.loaders.base_loader import BaseLoader


class PPTXLoader(BaseLoader):
    """
    Loads PowerPoint (.pptx) presentations.

    Each slide becomes one LangChain Document.
    """

    def load(self, file_path: Path):

        presentation = Presentation(file_path)

        documents = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):

            texts = []

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    if shape.text.strip():

                        texts.append(shape.text)

            if texts:

                documents.append(
                    LangChainDocument(
                        page_content="\n".join(texts),
                        metadata={
                            "source": file_path.name,
                            "page": slide_number,
                            "file_type": "pptx",
                        },
                    )
                )

        return documents